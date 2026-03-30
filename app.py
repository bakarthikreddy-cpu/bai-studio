import streamlit as st
from groq import Groq
from datetime import datetime
import fitz
from docx import Document
from pptx import Presentation
import os
import json
import csv
import io as io_module
import base64
import openpyxl
import time

# ---------- PAGE CONFIG ----------
st.set_page_config(
    page_title="BAi Studio",
    page_icon="🎯",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ---------- GLOBAL CSS (DARK THEME + ANIMATIONS + RESPONSIVE) ----------
st.markdown("""
<style>
    /* Global layout */
    .stApp {
        background: radial-gradient(circle at top, #111827 0, #020617 40%, #020617 100%);
        color: #e5e7eb;
        font-family: "Inter", system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }
    .block-container {
        padding: 1.5rem 2.4rem 3rem 2.4rem;
        max-width: 1160px;
    }

    /* Hide default Streamlit header */
    header[data-testid="stHeader"] {
        background: transparent;
    }

    /* Animated hero background */
    .hero-wrapper {
        position: relative;
        border-radius: 24px;
        padding: 2.6rem 2.4rem;
        overflow: hidden;
        border: 1px solid rgba(37, 99, 235, 0.4);
        background: radial-gradient(circle at 0% 0%, rgba(37, 99, 235, 0.35), transparent 55%),
                    radial-gradient(circle at 100% 100%, rgba(6, 182, 212, 0.28), transparent 55%),
                    #020617;
        box-shadow:
            0 18px 40px rgba(15, 23, 42, 0.9),
            0 0 120px rgba(37, 99, 235, 0.25);
    }
    .hero-orbit {
        position: absolute;
        inset: -40%;
        background:
            radial-gradient(circle at 10% 20%, rgba(56,189,248,0.22) 0, transparent 60%),
            radial-gradient(circle at 80% 80%, rgba(129,140,248,0.18) 0, transparent 60%);
        mix-blend-mode: screen;
        opacity: 0.7;
        filter: blur(4px);
        animation: hero-pulse 14s ease-in-out infinite alternate;
        pointer-events: none;
    }
    @keyframes hero-pulse {
        0%   { transform: translate3d(0, 0, 0) scale(1); opacity: 0.7; }
        50%  { transform: translate3d(-12px, -4px, 0) scale(1.03); opacity: 0.9; }
        100% { transform: translate3d(10px, 8px, 0) scale(1.05); opacity: 0.75; }
    }

    .hero-logo {
        font-size: 2.5rem;
        font-weight: 900;
        letter-spacing: 0.06em;
        text-transform: uppercase;
        background: linear-gradient(120deg, #60a5fa, #22d3ee, #a855f7);
        -webkit-background-clip: text;
        background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.2rem;
    }
    .hero-tagline {
        color: #9ca3af;
        font-size: 0.98rem;
        max-width: 520px;
    }

    /* Typing indicator */
    .typing-pill {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        padding: 4px 10px;
        border-radius: 999px;
        background: rgba(15,23,42,0.9);
        border: 1px solid rgba(148,163,184,0.4);
        color: #9ca3af;
        font-size: 0.75rem;
        margin-top: 0.8rem;
    }
    .typing-dot {
        width: 5px;
        height: 5px;
        border-radius: 999px;
        background: #60a5fa;
        animation: typing-bounce 1.2s infinite ease-in-out;
    }
    .typing-dot:nth-child(2) { animation-delay: 0.12s; }
    .typing-dot:nth-child(3) { animation-delay: 0.24s; }
    @keyframes typing-bounce {
        0%, 80%, 100%  { transform: translateY(0); opacity: 0.4; }
        40%            { transform: translateY(-3px); opacity: 1; }
    }

    /* Badges / hero chips */
    .badge {
        display: inline-flex;
        align-items: center;
        gap: 6px;
        background: rgba(15,23,42,0.9);
        color: #e5e7eb;
        border-radius: 999px;
        padding: 4px 12px;
        font-size: 0.72rem;
        font-weight: 600;
        margin-right: 8px;
        margin-top: 8px;
        border: 1px solid rgba(37,99,235,0.5);
        box-shadow: 0 0 14px rgba(37, 99, 235, 0.25);
        backdrop-filter: blur(10px);
    }
    .badge-dot {
        width: 6px;
        height: 6px;
        border-radius: 999px;
        background: #22c55e;
        box-shadow: 0 0 8px rgba(34, 197, 94, 0.9);
    }

    /* Sidebar branding */
    div[data-testid="stSidebar"] {
        background: #020617;
        border-right: 1px solid rgba(15,23,42,0.9);
        box-shadow: 8px 0 24px rgba(15,23,42,0.8);
    }
    .sidebar-brand {
        font-size: 1.4rem;
        font-weight: 800;
        letter-spacing: 0.08em;
        text-transform: uppercase;
        background: linear-gradient(135deg, #60a5fa, #22d3ee);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
    }
    .sidebar-subtitle {
        color: #64748b;
        font-size: 0.75rem;
        margin-bottom: 0.6rem;
        text-transform: uppercase;
        letter-spacing: 0.12em;
    }
    .sidebar-section-label {
        color: #94a3b8;
        font-size: 0.78rem;
        font-weight: 600;
        margin-bottom: 0.3rem;
        margin-top: 0.9rem;
        text-transform: uppercase;
        letter-spacing: 0.1em;
    }

    /* Metric cards */
    .metric-card {
        background: radial-gradient(circle at top left, rgba(37,99,235,0.52), transparent 55%),
                    #020617;
        border-radius: 14px;
        padding: 0.9rem 1rem;
        border: 1px solid rgba(30,64,175,0.7);
        text-align: left;
        position: relative;
        overflow: hidden;
    }
    .metric-card::after {
        content: "";
        position: absolute;
        inset: -40%;
        background: radial-gradient(circle at 10% 0%, rgba(56,189,248,0.18), transparent 55%);
        opacity: 0.2;
        pointer-events: none;
    }
    .metric-value {
        font-size: 1.7rem;
        font-weight: 800;
        background: linear-gradient(135deg, #60a5fa, #22d3ee);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        margin-bottom: 0.15rem;
        position: relative;
        z-index: 1;
    }
    .metric-label {
        font-size: 0.7rem;
        color: #9ca3af;
        text-transform: uppercase;
        letter-spacing: 0.12em;
        font-weight: 600;
        position: relative;
        z-index: 1;
    }

    /* Tabs refinement */
    button[data-baseweb="tab"] {
        background: transparent !important;
        border-radius: 999px !important;
        padding: 0.45rem 1.2rem !important;
        margin-right: 0.35rem !important;
        color: #9ca3af !important;
        border: 1px solid transparent !important;
    }
    button[data-baseweb="tab"][aria-selected="true"] {
        background: linear-gradient(135deg, rgba(37,99,235,0.15), rgba(6,182,212,0.12)) !important;
        border-color: rgba(37,99,235,0.7) !important;
        color: #e5e7eb !important;
        box-shadow: 0 10px 28px rgba(15,23,42,0.7);
    }

    /* Inputs & widgets */
    .stTextArea textarea, .stTextInput input {
        background: #020617 !important;
        border-radius: 12px !important;
        border: 1px solid rgba(51,65,85,0.9) !important;
        color: #e5e7eb !important;
        font-size: 0.92rem !important;
    }
    .stSelectbox select, .stMultiSelect div[data-baseweb="select"] {
        background: #020617 !important;
        border-radius: 10px !important;
        border: 1px solid rgba(51,65,85,0.9) !important;
        color: #e5e7eb !important;
    }
    .stSlider > div > div > div:nth-child(2) > div {
        background: linear-gradient(90deg, #60a5fa, #22d3ee) !important;
    }
    .stSlider > div > div > div:nth-child(2) > div > div {
        box-shadow: 0 0 0 4px rgba(37,99,235,0.35);
    }

    /* Primary button */
    .stButton > button {
        background: radial-gradient(circle at 0 0, #38bdf8, #2563eb);
        color: #f9fafb;
        border-radius: 999px;
        border: none;
        padding: 0.7rem 1.4rem;
        font-size: 0.95rem;
        font-weight: 700;
        width: 100%;
        box-shadow: 0 12px 30px rgba(37,99,235,0.55);
        transition: all 0.16s ease-out;
    }
    .stButton > button:hover {
        transform: translateY(-1px) scale(1.01);
        box-shadow: 0 16px 40px rgba(37,99,235,0.7);
        filter: brightness(1.05);
    }

    /* Result card */
    .result-box {
        background: radial-gradient(circle at top left, rgba(37,99,235,0.18), transparent 55%),
                    rgba(15,23,42,0.98);
        border-left: 4px solid #38bdf8;
        border-radius: 0 16px 16px 16px;
        padding: 1.4rem 1.6rem;
        color: #e5e7eb;
        line-height: 1.8;
        font-size: 0.95rem;
        box-shadow: 0 20px 40px rgba(15,23,42,0.85);
        position: relative;
        overflow: hidden;
    }
    .result-box::before {
        content: "";
        position: absolute;
        inset: -40%;
        background: radial-gradient(circle at 20% 0%, rgba(56,189,248,0.18), transparent 55%);
        opacity: 0.4;
        pointer-events: none;
    }

    /* File preview */
    .file-preview {
        background: rgba(15,23,42,0.96);
        border-radius: 10px;
        padding: 0.6rem 0.9rem;
        border: 1px solid rgba(51,65,85,0.9);
        margin-bottom: 0.45rem;
        font-size: 0.8rem;
        color: #9ca3af;
        display: flex;
        justify-content: space-between;
        align-items: center;
    }

    /* Search result card */
    .search-result {
        background: rgba(15,23,42,0.96);
        border-left: 4px solid #2563eb;
        border-radius: 0 12px 12px 0;
        padding: 0.9rem 1.1rem;
        margin-bottom: 0.7rem;
        font-size: 0.84rem;
        color: #9ca3af;
        border-top: 1px solid rgba(37,99,235,0.5);
        border-bottom: 1px solid rgba(37,99,235,0.5);
    }

    /* Footer box */
    .footer-box {
        background: radial-gradient(circle at top, rgba(15,23,42,0.9), #020617);
        border-radius: 14px;
        padding: 1.1rem 1.6rem;
        border: 1px solid rgba(30,64,175,0.7);
        text-align: center;
        margin-top: 2.3rem;
        color: #9ca3af;
        font-size: 0.82rem;
    }
    .footer-box a {
        color: #38bdf8;
        text-decoration: none;
        font-weight: 600;
    }

    /* Copy button row */
    .copy-row {
        display: flex;
        justify-content: flex-end;
        gap: 0.5rem;
        margin-top: 0.5rem;
        margin-bottom: 0.4rem;
    }
    .copy-chip {
        font-size: 0.74rem;
        padding: 0.28rem 0.9rem;
        border-radius: 999px;
        border: 1px solid rgba(55,65,81,0.9);
        background: rgba(15,23,42,0.98);
        color: #9ca3af;
    }

    /* Responsive tweaks */
    @media (max-width: 768px) {
        .block-container {
            padding: 1rem 1.2rem 2.3rem 1.2rem;
        }
        .hero-logo {
            font-size: 2rem;
        }
        .hero-tagline {
            font-size: 0.9rem;
        }
        .hero-wrapper {
            padding: 1.9rem 1.6rem;
        }
    }
</style>
""", unsafe_allow_html=True)

# ---------- SESSION STATE ----------
if "history" not in st.session_state:
    st.session_state.history = []
if "run_count" not in st.session_state:
    st.session_state.run_count = 0
if "last_result" not in st.session_state:
    st.session_state.last_result = ""
if "ba_prefill" not in st.session_state:
    st.session_state.ba_prefill = ""

# ---------- HELPERS ----------
def extract_pdf(file):
    text = ""
    doc = fitz.open(stream=file.read(), filetype="pdf")
    for page in doc:
        text += page.get_text()
    return text.strip()

def extract_docx(file):
    doc = Document(file)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extract_pptx(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_image(file):
    try:
        groq_key = st.secrets.get("GROQ_API_KEY", os.environ.get("GROQ_API_KEY", ""))
        client = Groq(api_key=groq_key)
        image_data = base64.b64encode(file.read()).decode("utf-8")
        file_name = file.name.lower()
        media_type = "image/png" if file_name.endswith(".png") else "image/jpeg"
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Extract all text and describe all content from this image in detail."},
                    {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{image_data}"}}
                ]
            }],
            max_tokens=2048
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"Image analysis unavailable: {str(e)}"

def extract_json(file):
    try:
        data = json.load(file)
        return json.dumps(data, indent=2)
    except Exception as e:
        return f"Error reading JSON: {str(e)}"

def extract_excel(file):
    try:
        wb = openpyxl.load_workbook(file)
        text = ""
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text += f"\n--- Sheet: {sheet} ---\n"
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join([str(c) if c is not None else "" for c in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text.strip()
    except Exception as e:
        return f"Error reading Excel: {str(e)}"

def extract_csv(file):
    try:
        content = file.read().decode("utf-8")
        reader = csv.reader(io_module.StringIO(content))
        rows = ["\t".join(row) for row in reader]
        return "\n".join(rows)
    except Exception as e:
        return f"Error reading CSV: {str(e)}"

def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        return extract_pdf(uploaded_file)
    elif name.endswith((".docx", ".doc")):
        return extract_docx(uploaded_file)
    elif name.endswith(".pptx"):
        return extract_pptx(uploaded_file)
    elif name.endswith((".png", ".jpg", ".jpeg", ".tiff", ".bmp")):
        return extract_image(uploaded_file)
    elif name.endswith(".txt"):
        return uploaded_file.read().decode("utf-8")
    elif name.endswith(".json"):
        return extract_json(uploaded_file)
    elif name.endswith((".xlsx", ".xls")):
        return extract_excel(uploaded_file)
    elif name.endswith(".csv"):
        return extract_csv(uploaded_file)
    return None

def universal_uploader(key):
    uploaded = st.file_uploader(
        "📎 Upload documents for context (optional)",
        type=[
            "pdf","docx","doc","pptx","png","jpg","jpeg","bmp",
            "tiff","txt","json","xlsx","xls","csv"
        ],
        accept_multiple_files=True,
        key=f"uploader_{key}"
    )
    context = ""
    if uploaded:
        st.markdown("**📂 Files loaded into context:**")
        for f in uploaded:
            size_kb = round(f.size / 1024, 1)
            st.markdown(
                f"<div class='file-preview'>"
                f"<span>📄 {f.name}</span>"
                f"<span>{size_kb} KB</span>"
                f"</div>",
                unsafe_allow_html=True
            )
            extracted = extract_text_from_file(f)
            if extracted:
                context += f"\n\n--- {f.name} ---\n{extracted[:2000]}"
    return context

def web_search(query):
    try:
        from duckduckgo_search import DDGS
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=3):
                results.append(f"- {r['title']}: {r['body']}")
        return "\n".join(results)
    except Exception as e:
        return f"Search unavailable: {str(e)}"

def run_crew(scenario, role, num_agents, model, temperature=0.7, progress_placeholder=None):
    groq_key = st.secrets.get("GROQ_API_KEY", os.environ.get("GROQ_API_KEY", ""))
    model_map = {
        "groq/llama-3.3-70b-versatile": "llama-3.3-70b-versatile",
        "groq/llama-3.1-8b-instant": "llama-3.1-8b-instant",
        "groq/gemma2-9b-it": "gemma2-9b-it",
        "groq/compound-beta": "compound-beta"
    }
    groq_model = model_map.get(model, "llama-3.3-70b-versatile")
    personas = [
        (
            "Senior Business Analyst",
            "You are BAi the Ultimate Business Analyst Agent. You are a CBAP-certified Senior BA with 15 years "
            "experience in BABOK v3, Agile, SAFe, Waterfall, SAP, Salesforce and Kronos WFM. Always structure your "
            "output with: Executive Summary, Key Findings, Recommendations, and Next Steps. Be thorough, precise and indispensable."
        ),
        (
            "Stakeholder Advocate",
            "You represent all stakeholder groups. Challenge assumptions and identify gaps, conflicts and missing "
            "requirements from a stakeholder perspective. Always highlight risks and open questions."
        ),
        (
            "Solution Designer",
            "You synthesize BA findings into clear actionable solution recommendations. Focus on feasibility, "
            "business value and implementation approach. Provide a prioritized action plan."
        )
    ]
    client = Groq(api_key=groq_key)
    combined_result = ""
    for i in range(num_agents):
        if progress_placeholder is not None:
            progress_placeholder.progress((i + 1) / max(num_agents, 1))
        agent_role, backstory = personas[i]
        messages = [
            {"role": "system", "content": f"You are a {agent_role}. {backstory}"},
            {
                "role": "user",
                "content": (
                    f"Analyze this scenario: {scenario}\n\n"
                    f"Provide a structured analysis with: Executive Summary, Key Findings, Recommendations, and Next Steps."
                )
            }
        ]
        try:
            response = client.chat.completions.create(
                model=groq_model,
                messages=messages,
                temperature=temperature,
                max_tokens=2048
            )
            agent_output = response.choices[0].message.content
            if num_agents > 1:
                combined_result += f"\n\n---\n### Agent: {agent_role}\n{agent_output}"
            else:
                combined_result = agent_output
        except Exception as e:
            combined_result += f"\n\nError from {agent_role}: {str(e)}"
    if progress_placeholder is not None:
        progress_placeholder.empty()
    return combined_result

def typewriter_effect(text, speed=120):
    """Lightweight typing effect for short strings."""
    container = st.empty()
    current = ""
    for ch in text:
        current += ch
        container.markdown(f"**{current}**")
        time.sleep(1 / speed)

# ---------- SIDEBAR ----------
with st.sidebar:
    st.markdown("<div class='sidebar-brand'>BAi Studio</div>", unsafe_allow_html=True)
    st.markdown("<div class='sidebar-subtitle'>Business Analysis Intelligence</div>", unsafe_allow_html=True)
    st.markdown("---")

    st.markdown("<div class='sidebar-section-label'>AI MODEL</div>", unsafe_allow_html=True)
    model_choice = st.selectbox(
        "Model",
        [
            "groq/llama-3.3-70b-versatile",
            "groq/llama-3.1-8b-instant",
            "groq/gemma2-9b-it",
            "groq/compound-beta"
        ],
        label_visibility="collapsed"
    )

    st.markdown("<div class='sidebar-section-label'>AGENTS</div>", unsafe_allow_html=True)
    num_agents = st.slider("Agents", 1, 3, 1, label_visibility="collapsed")
    agent_names = [
        "Senior BA only",
        "Senior BA + Stakeholder Advocate",
        "Full Crew (BA + Advocate + Designer)"
    ]
    st.markdown(
        f"<div style='color:#38bdf8;font-size:0.76rem;font-weight:600'>Active: {agent_names[num_agents-1]}</div>",
        unsafe_allow_html=True
    )

    st.markdown("<div class='sidebar-section-label'>DEPTH</div>", unsafe_allow_html=True)
    depth = st.select_slider(
        "Depth",
        options=["Quick", "Standard", "Deep"],
        value="Standard",
        label_visibility="collapsed"
    )

    st.markdown("<div class='sidebar-section-label'>CREATIVITY</div>", unsafe_allow_html=True)
    temperature = st.slider("Temperature", 0.1, 1.0, 0.7, 0.1, label_visibility="collapsed")

    st.markdown("<div class='sidebar-section-label'>WEB SEARCH</div>", unsafe_allow_html=True)
    use_search = st.toggle("Search web before analysis", value=False)

    st.markdown("---")

    col1, col2 = st.columns(2)
    with col1:
        st.markdown(
            f"<div class='metric-card'><div class='metric-value'>{st.session_state.run_count}</div>"
            f"<div class='metric-label'>Runs</div></div>",
            unsafe_allow_html=True
        )
    with col2:
        st.markdown(
            f"<div class='metric-card'><div class='metric-value'>{len(st.session_state.history)}</div>"
            f"<div class='metric-label'>Saved</div></div>",
            unsafe_allow_html=True
        )

    if st.session_state.history:
        st.markdown("---")
        st.markdown(
            "<div class='sidebar-section-label'>Recent runs</div>",
            unsafe_allow_html=True
        )
        for item in reversed(st.session_state.history[-4:]):
            st.markdown(
                f"<div style='background:#020617;border-radius:9px;padding:0.55rem;"
                f"margin-bottom:0.38rem;border:1px solid rgba(30,64,175,0.7);"
                f"font-size:0.8rem;color:#9ca3af'>"
                f"<strong>{item['type']}</strong><br>"
                f"<span style='color:#6b7280;font-size:0.72rem'>{item['time']}</span>"
                f"</div>",
                unsafe_allow_html=True
            )

    st.markdown("---")
    st.markdown(
        "<a href='https://www.linkedin.com/in/karthik-reddy-t-666334232/' "
        "target='_blank' style='color:#38bdf8;text-decoration:none;font-size:0.82rem;font-weight:600'>"
        "Karthik Reddy T · LinkedIn</a>",
        unsafe_allow_html=True
    )

# ---------- HERO ----------
st.markdown(
    """
    <div class='hero-wrapper'>
        <div class='hero-orbit'></div>
        <div style='position:relative;z-index:1;display:flex;flex-wrap:wrap;gap:2rem;align-items:flex-start;justify-content:space-between'>
            <div style='flex:1 1 280px;min-width:260px'>
                <div class='hero-logo'>BAi Studio</div>
                <div class='hero-tagline'>
                    Your AI-powered workspace for turning messy documents and stakeholder noise
                    into clear, decision-ready business analysis.
                </div>
                <div class='typing-pill'>
                    <span class='typing-dot'></span>
                    <span class='typing-dot'></span>
                    <span class='typing-dot'></span>
                    <span>BAi is thinking in BABOK...</span>
                </div>
                <br><br>
                <span class='badge'><span class='badge-dot'></span>BABOK Aligned</span>
                <span class='badge'>CBAP Ready</span>
                <span class='badge'>Smart Document Analysis</span>
            </div>
            <div style='flex:0 0 260px;min-width:220px'>
                <div style='background:rgba(15,23,42,0.95);border-radius:18px;padding:1rem 1.2rem;border:1px solid rgba(30,64,175,0.85);box-shadow:0 18px 40px rgba(15,23,42,0.9)'>
                    <div style='font-size:0.8rem;color:#9ca3af;margin-bottom:0.4rem;font-weight:600;text-transform:uppercase;letter-spacing:0.12em'>
                        Quick Start
                    </div>
                    <ul style='margin:0;padding-left:1.1rem;font-size:0.84rem;color:#d1d5db;line-height:1.7'>
                        <li>Describe a messy BA scenario</li>
                        <li>Upload BRDs, SoWs, notes, or images</li>
                        <li>Pick depth & agents in the sidebar</li>
                        <li>Get CBAP-style, ready-to-ship output</li>
                    </ul>
                </div>
            </div>
        </div>
    </div>
    """,
    unsafe_allow_html=True
)

# Optional tiny typewriter hint under hero
with st.container():
    typewriter_effect("Built for real-world BA work, not toy demos.", speed=90)

# ---------- TABS ----------
tab1, tab2, tab3, tab4 = st.tabs([
    "Requirements Analysis",
    "Document Analyzer",
    "BA Toolkit",
    "Session History"
])

# ---------- TAB 1: Requirements Analysis ----------
with tab1:
    st.markdown("#### Requirements & Business Analysis")
    st.markdown(
        "<div style='color:#9ca3af;font-size:0.9rem;margin-bottom:1rem'>"
        "Describe your BA scenario, problem statement, or requirements challenge. BAi will respond "
        "with CBAP-style structure and language."
        "</div>",
        unsafe_allow_html=True
    )

    babok_area = st.selectbox(
        "BABOK Knowledge Area",
        [
            "Business Analysis Planning & Monitoring",
            "Elicitation & Collaboration",
            "Requirements Life Cycle Management",
            "Strategy Analysis",
            "Requirements Analysis & Design Definition",
            "Solution Evaluation"
        ]
    )

    scenario = st.text_area(
        "Scenario",
        placeholder="E.g. A company wants to implement a new CRM system. Identify stakeholders, "
                    "elicit requirements, and define success criteria.",
        height=140,
        label_visibility="collapsed"
    )

    doc_context_tab1 = universal_uploader("tab1")

    output_format = st.selectbox(
        "Output Format",
        [
            "Full Structured Report",
            "Executive Summary Only",
            "Bullet Points",
            "BA Deliverable Format"
        ]
    )

    run_req = st.button("Run BA Analysis", type="primary", key="req_run")

    if run_req:
        if scenario.strip():
            progress = st.progress(0)
            status = st.empty()
            status.info("Preparing scenario and context...")

            final_scenario = (
                f"BABOK Area: {babok_area}\n"
                f"Output Format: {output_format}\n"
                f"Depth: {depth}\n"
                f"Scenario: {scenario}"
            )
            if doc_context_tab1:
                final_scenario += f"\n\nUploaded Documents:\n{doc_context_tab1}"

            if use_search:
                with st.spinner("Searching web for latest context..."):
                    search_data = web_search(scenario[:120])
                    if search_data:
                        st.markdown("<div class='search-result'>" + search_data + "</div>", unsafe_allow_html=True)
                        final_scenario += "\n\nRecent web data:\n" + search_data

            status.info("BA agents analyzing scenario...")
            result_str = run_crew(
                final_scenario,
                babok_area,
                num_agents,
                model_choice,
                temperature,
                progress_placeholder=progress
            )

            st.session_state.run_count += 1
            st.session_state.history.append({
                "type": babok_area[:40],
                "tab": "Requirements",
                "scenario": scenario[:80] + ("..." if len(scenario) > 80 else ""),
                "time": datetime.now().strftime("%b %d, %I:%M %p"),
                "result": result_str,
                "model": model_choice
            })
            st.session_state.last_result = result_str

            status.success("BA Analysis Complete!")
            st.balloons()  # confetti-style celebration

            st.markdown("<div class='copy-row'>", unsafe_allow_html=True)
            copy_col1, copy_col2 = st.columns([1, 1])
            with copy_col1:
                if st.button("Copy to clipboard (browser)", key="req_copy"):
                    st.write("Use your browser's copy from the result area.")
            with copy_col2:
                st.download_button(
                    "Download Report",
                    data=result_str,
                    file_name=f"ba_analysis_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain",
                    key="req_download"
                )
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("<div class='result-box'>" + result_str + "</div>", unsafe_allow_html=True)
        else:
            st.warning("Please describe your scenario first!")

# ---------- TAB 2: Document Analyzer ----------
with tab2:
    st.markdown("#### Document Analyzer")
    st.markdown(
        "<div style='color:#9ca3af;font-size:0.9rem;margin-bottom:1rem'>"
        "Upload BRDs, SOWs, contracts, meeting notes, process docs, images, Excel, JSON — "
        "any business document that needs BA eyes."
        "</div>",
        unsafe_allow_html=True
    )

    uploaded_files = st.file_uploader(
        "Upload Documents",
        type=["pdf","docx","doc","pptx","png","jpg","jpeg","bmp","tiff","txt","json","xlsx","xls","csv"],
        accept_multiple_files=True,
        label_visibility="collapsed"
    )

    analysis_type = st.selectbox(
        "What do you want to extract?",
        [
            "Extract & Summarize Requirements",
            "Identify Stakeholders",
            "Flag Risks & Assumptions",
            "List Action Items & Decisions",
            "Find Missing Requirements",
            "Full BA Document Review"
        ]
    )

    doc_question = st.text_area(
        "Additional Instructions (optional)",
        placeholder="E.g. Focus on functional requirements only. Ignore section 3.",
        height=80
    )

    if uploaded_files:
        st.markdown("**Uploaded Files:**")
        for f in uploaded_files:
            size_kb = round(f.size / 1024, 1)
            st.markdown(
                f"<div class='file-preview'><span>{f.name}</span><span>{size_kb} KB</span></div>",
                unsafe_allow_html=True
            )

    run_doc = st.button("Analyze Documents", type="primary", key="doc_run")

    if run_doc:
        if not uploaded_files:
            st.warning("Please upload at least one file!")
        else:
            progress = st.progress(0)
            status = st.empty()
            status.info("Extracting text from documents...")

            all_text = ""
            for idx, f in enumerate(uploaded_files):
                extracted = extract_text_from_file(f)
                if extracted:
                    all_text += f"\n\n--- {f.name} ---\n{extracted}"
                progress.progress((idx + 1) / max(len(uploaded_files), 1))

            if all_text.strip():
                status.info("Senior BA analyzing documents...")
                full_scenario = (
                    f"Task: {analysis_type}\n"
                    f"Additional Instructions: {doc_question}\n"
                    f"Depth: {depth}\n"
                    f"Document Content:\n{all_text[:4000]}"
                )
                result_str = run_crew(
                    full_scenario,
                    "Document Analysis Specialist",
                    1,
                    model_choice,
                    temperature,
                    progress_placeholder=progress
                )

                st.session_state.run_count += 1
                st.session_state.history.append({
                    "type": "Document Analysis",
                    "tab": "Document",
                    "scenario": analysis_type,
                    "time": datetime.now().strftime("%b %d, %I:%M %p"),
                    "result": result_str,
                    "model": model_choice
                })
                st.session_state.last_result = result_str

                status.success("Document Analysis Complete!")
                st.balloons()

                st.markdown("<div class='copy-row'>", unsafe_allow_html=True)
                copy_col1, copy_col2 = st.columns([1, 1])
                with copy_col1:
                    if st.button("Copy to clipboard (browser)", key="doc_copy"):
                        st.write("Use your browser's copy from the result area.")
                with copy_col2:
                    st.download_button(
                        "Download Analysis",
                        data=result_str,
                        file_name=f"doc_analysis_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain",
                        key="doc_download"
                    )
                st.markdown("</div>", unsafe_allow_html=True)

                st.markdown("<div class='result-box'>" + result_str + "</div>", unsafe_allow_html=True)
            else:
                status.empty()
                st.error("Could not extract text from uploaded files.")

# ---------- TAB 3: BA Toolkit ----------
with tab3:
    st.markdown("#### BA Toolkit")
    st.markdown(
        "<div style='color:#9ca3af;font-size:0.9rem;margin-bottom:1.5rem'>"
        "One-click BA deliverable generators – all BABOK-aligned and CBAP-inspired."
        "</div>",
        unsafe_allow_html=True
    )

    ba_templates = {
        "User Stories": "Write 5 detailed user stories in INVEST format for a digital transformation project in Ontario, Canada. Include acceptance criteria for each.",
        "Gap Analysis": "Perform a detailed gap analysis for a business switching from manual processes to an ERP system in Ontario. Include AS-IS state, TO-BE state, and gap recommendations.",
        "Stakeholder Register": "Create a complete stakeholder register for implementing a new CRM system at a mid-size business in Ontario. Include interest, influence, and engagement strategy for each stakeholder.",
        "BRD Generator": "Write a Business Requirements Document section for automating payroll processing using Kronos WFM. Include business objectives, scope, assumptions, and functional requirements.",
        "CBAP Practice": "Generate 10 CBAP exam-style scenario questions for the Requirements Life Cycle Management knowledge area with detailed answers and BABOK references.",
        "Process Flow": "Document the AS-IS and TO-BE process flow for a manual invoice approval process being automated with SAP. Include swimlanes, decision points, and improvement recommendations.",
        "Risk Register": "Create a BA risk register for an ERP implementation project. Include risk description, probability, impact, mitigation strategy, and owner for each risk.",
        "Use Case": "Write 3 detailed use cases for an online customer portal. Include actors, preconditions, main flow, alternate flow, and postconditions.",
        "Meeting Notes": "Analyze these meeting notes and extract: key decisions made, action items with owners, open issues, risks identified, and next steps.",
        "Business Case": "Write a mini business case for implementing an AI-powered document management system at a mid-size Canadian company. Include problem statement, options analysis, costs, benefits, and recommendation."
    }

    cols = st.columns(2)
    for idx, (label, prompt) in enumerate(ba_templates.items()):
        with cols[idx % 2]:
            if st.button(label, key=f"ba_{idx}"):
                st.session_state.ba_prefill = prompt
                st.rerun()

    st.markdown("---")

    ba_scenario = st.text_area(
        "BA Task",
        value=st.session_state.ba_prefill,
        height=130,
        key="ba_input"
    )

    doc_context_tab3 = universal_uploader("tab3")

    run_ba = st.button("Generate BA Deliverable", type="primary", key="ba_run")

    if run_ba:
        if ba_scenario.strip():
            progress = st.progress(0)
            status = st.empty()
            status.info("CBAP-certified BA agent generating deliverable...")

            full_ba = (
                f"Depth: {depth}\n"
                f"BA Task:\n{ba_scenario}"
            )
            if doc_context_tab3:
                full_ba += f"\n\nUploaded Documents:\n{doc_context_tab3}"

            result_str = run_crew(
                full_ba,
                "Senior Business Analyst",
                1,
                model_choice,
                temperature,
                progress_placeholder=progress
            )

            st.session_state.run_count += 1
            st.session_state.history.append({
                "type": "BA Toolkit",
                "tab": "BA Toolkit",
                "scenario": ba_scenario[:80],
                "time": datetime.now().strftime("%b %d, %I:%M %p"),
                "result": result_str,
                "model": model_choice
            })
            st.session_state.last_result = result_str

            status.success("BA Deliverable Ready!")
            st.balloons()

            st.markdown("<div class='copy-row'>", unsafe_allow_html=True)
            copy_col1, copy_col2 = st.columns([1, 1])
            with copy_col1:
                if st.button("Copy to clipboard (browser)", key="ba_copy"):
                    st.write("Use your browser's copy from the result area.")
            with copy_col2:
                st.download_button(
                    "Download Deliverable",
                    data=result_str,
                    file_name=f"ba_deliverable_{datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain",
                    key="ba_download"
                )
            st.markdown("</div>", unsafe_allow_html=True)

            st.markdown("<div class='result-box'>" + result_str + "</div>", unsafe_allow_html=True)
        else:
            st.warning("Please select a template or enter a BA task!")

# ---------- TAB 4: Session History ----------
with tab4:
    st.markdown("#### Session History")
    if not st.session_state.history:
        st.info("No analyses run yet. Start with Requirements Analysis or BA Toolkit!")
    else:
        st.markdown(f"**Total analyses this session: {st.session_state.run_count}**")
        st.markdown("---")
        for i, item in enumerate(reversed(st.session_state.history)):
            run_num = len(st.session_state.history) - i
            with st.expander(f"Run #{run_num} · {item['type']} · {item['time']}"):
                st.markdown(f"**Tab:** {item['tab']} | **Model:** {item.get('model','N/A')}")
                st.markdown(f"**Scenario:** {item['scenario']}")
                st.markdown("**Result:**")
                st.markdown("<div class='result-box'>" + item["result"] + "</div>", unsafe_allow_html=True)
                st.download_button(
                    f"Download Run #{run_num}",
                    data=item["result"],
                    file_name=f"run_{run_num}.txt",
                    mime="text/plain",
                    key=f"dl_{i}"
                )

        if st.button("Clear Session History", key="clear_history"):
            st.session_state.history = []
            st.session_state.run_count = 0
            st.session_state.last_result = ""
            st.rerun()

# ---------- FOOTER ----------
st.markdown("""
<div class='footer-box'>
    <strong>BAi Studio</strong> · Your AI-Powered Business Analysis Workspace<br><br>
    Built by <a href='https://www.linkedin.com/in/karthik-reddy-t-666334232/' target='_blank'>Karthik Reddy T</a> · Business Analyst · CBAP Candidate
</div>
""", unsafe_allow_html=True)